from celery import shared_task
# from PyPDF2 import PdfReader
from pptx import Presentation
import pandas as pd
from zipfile import ZipFile 
import docx
from app.tools import getWhatonImage_celery, get_ext_and_mime, gemini_upload_and_chat
from io import BytesIO

@shared_task
def extract_text_from_file(file_bytes):
    file_ext = get_ext_and_mime(file_bytes)
    print(f"File extension: {file_ext}")
    
    if file_ext == 'zip':
        extracted_files = extract_zip(file_bytes)
        text = ""
        for _, content in extracted_files.items():
            text += extract_text_from_file(content)
        return text
    
    elif file_ext == 'application/pdf':
        return gemini_upload_and_chat(bytes=file_bytes, mime=file_ext)
    
    elif file_ext == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
        return get_xls_text(file_bytes)
    
    elif file_ext == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return get_pptx_text(file_bytes)
    
    elif file_ext == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return getText(file_bytes)
    
    elif file_ext in ['image/jpeg', 'image/png']:
        return get_text_from_image(file_bytes)
    
    else:
        return "Unsupported file type"

def get_xls_text(xls_Bytes):
    text = ""
    xls = pd.ExcelFile(BytesIO(xls_Bytes))
    for sheet in xls.sheet_names:
        df = pd.read_excel(xls, sheet)
        text += df.to_string()
    return text

def extract_zip(input_zip):
    input_zip=ZipFile(input_zip)
    return {name: input_zip.read(name) for name in input_zip.namelist()}

def get_text_from_image(image_files):
    return getWhatonImage_celery.delay(image_files).get()

# def get_pdf_text(pdf_files):
#     text = ""
#     for pdf in pdf_files:
#         pdf_reader = PdfReader(pdf)
#         for page in pdf_reader.pages:
#             text += page.extract_text()
#     return text

def getText(fileBytes):
    doc = docx.Document(BytesIO(fileBytes))
    fullText = [para.text for para in doc.paragraphs]
    return '\n'.join(fullText)

def get_pptx_text(pptxBytes):
    text = ""
    prs = Presentation(BytesIO(pptxBytes))
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                fullText.append(shape.text)
    text += '\n'.join(fullText)
    return text


